#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Test data generator - Creates sample Japanese-Vietnamese bilingual files
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

# Sample bilingual data (Japanese - Vietnamese pairs)
# Construction/Safety related terms
TEST_DATA = [
    {
        "title_ja": "安全作業手順書",
        "title_vi": "Quy trinh lam viec an toan",
        "items": [
            ("ヘルメットを着用してください", "Hay doi mu bao hiem"),
            ("安全帯を必ず装着すること", "Phai deo day an toan"),
            ("足場の点検を行う", "Kiem tra gian giao"),
            ("作業前に危険箇所を確認", "Xac nhan vi tri nguy hiem truoc khi lam viec"),
        ]
    },
    {
        "title_ja": "建築用語集",
        "title_vi": "Tu dien thuat ngu xay dung",
        "items": [
            ("鉄筋コンクリート", "Be tong cot thep"),
            ("基礎工事", "Cong trinh nen mong"),
            ("配管工事", "Cong trinh duong ong"),
            ("電気設備", "Thiet bi dien"),
        ]
    },
    {
        "title_ja": "朝礼での注意事項",
        "title_vi": "Nhung diem can luu y trong buoi hop sang",
        "items": [
            ("今日の作業内容を確認します", "Xac nhan noi dung cong viec hom nay"),
            ("体調不良の方は申し出てください", "Neu ban khong khoe, hay bao cao"),
            ("熱中症に注意してください", "Hay chu y phong tranh say nang"),
            ("水分補給を忘れずに", "Dung quen uong nuoc"),
        ]
    },
    {
        "title_ja": "緊急時の対応",
        "title_vi": "Ung pho khi khan cap",
        "items": [
            ("事故が発生したら直ちに報告", "Bao cao ngay khi xay ra tai nan"),
            ("避難経路を確認してください", "Hay xac nhan duong thoat hiem"),
            ("消火器の場所を覚えておく", "Nho vi tri binh chua chay"),
            ("応急処置キットの使い方", "Cach su dung bo so cuu"),
        ]
    },
]


def create_bilingual_pptx(output_path: Path):
    """Create a bilingual PowerPoint file"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide layout
    title_layout = prs.slide_layouts[6]  # Blank layout

    for section in TEST_DATA:
        slide = prs.slides.add_slide(title_layout)

        # Add title (Japanese)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = section["title_ja"]
        title_para.font.size = Pt(32)
        title_para.font.bold = True

        # Add title (Vietnamese)
        title_vi_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.6))
        title_vi_frame = title_vi_box.text_frame
        title_vi_para = title_vi_frame.paragraphs[0]
        title_vi_para.text = section["title_vi"]
        title_vi_para.font.size = Pt(24)
        title_vi_para.font.italic = True

        # Add content items
        y_pos = 1.8
        for ja_text, vi_text in section["items"]:
            # Japanese text
            ja_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(6), Inches(0.5))
            ja_frame = ja_box.text_frame
            ja_para = ja_frame.paragraphs[0]
            ja_para.text = f"・{ja_text}"
            ja_para.font.size = Pt(18)

            # Vietnamese text
            vi_box = slide.shapes.add_textbox(Inches(6.5), Inches(y_pos), Inches(6), Inches(0.5))
            vi_frame = vi_box.text_frame
            vi_para = vi_frame.paragraphs[0]
            vi_para.text = f"・{vi_text}"
            vi_para.font.size = Pt(18)

            y_pos += 0.8

    prs.save(str(output_path))
    print(f"Created: {output_path}")


def create_test_csv(output_path: Path):
    """Create a CSV file with test pairs for direct testing"""
    import csv

    with open(output_path, 'w', encoding='utf-8-sig', newline='') as f:
        writer = csv.writer(f)
        writer.writerow(['ID', 'Japanese', 'Vietnamese', 'Category'])

        idx = 1
        for section in TEST_DATA:
            category = section["title_ja"]
            # Add title
            writer.writerow([idx, section["title_ja"], section["title_vi"], "Title"])
            idx += 1
            # Add items
            for ja_text, vi_text in section["items"]:
                writer.writerow([idx, ja_text, vi_text, category])
                idx += 1

    print(f"Created: {output_path}")


def main():
    output_dir = Path(__file__).parent / "test_files"
    output_dir.mkdir(exist_ok=True)

    # Create PowerPoint file
    pptx_path = output_dir / "safety_manual_ja_vi.pptx"
    create_bilingual_pptx(pptx_path)

    # Create CSV file (for reference/direct testing)
    csv_path = output_dir / "test_pairs.csv"
    create_test_csv(csv_path)

    print(f"\nTest files created in: {output_dir}")
    print(f"  - {pptx_path.name} (PowerPoint)")
    print(f"  - {csv_path.name} (CSV reference)")


if __name__ == "__main__":
    main()
