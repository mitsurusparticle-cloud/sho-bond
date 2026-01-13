"""
PowerPointファイルからテキストを抽出するモジュール
"""
from pathlib import Path
from typing import List, Tuple
from dataclasses import dataclass

from pptx import Presentation
from langdetect import detect, LangDetectException


@dataclass
class TextPair:
    """日本語・ベトナム語のテキストペア"""
    japanese: str
    vietnamese: str
    source_file: str
    slide_number: int
    confidence: float = 1.0


def detect_language(text: str) -> str:
    """
    テキストの言語を検出

    Args:
        text: 検出対象のテキスト

    Returns:
        言語コード (ja, vi, en, etc.) または "unknown"
    """
    if not text or len(text.strip()) < 3:
        return "unknown"

    try:
        return detect(text)
    except LangDetectException:
        return "unknown"


def extract_text_from_slide(slide) -> List[str]:
    """
    スライドから全テキストを抽出

    Args:
        slide: python-pptxのslideオブジェクト

    Returns:
        テキストのリスト
    """
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
        # テーブル内のテキストも抽出
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text.strip())
    return texts


def separate_by_language(texts: List[str]) -> Tuple[List[str], List[str]]:
    """
    テキストリストを日本語とベトナム語に分離

    Args:
        texts: テキストのリスト

    Returns:
        (日本語リスト, ベトナム語リスト)
    """
    japanese_texts = []
    vietnamese_texts = []

    for text in texts:
        lang = detect_language(text)
        if lang == "ja":
            japanese_texts.append(text)
        elif lang == "vi":
            vietnamese_texts.append(text)
        # その他の言語は無視

    return japanese_texts, vietnamese_texts


def extract_from_pptx(file_path: str | Path) -> List[TextPair]:
    """
    PowerPointファイルから日本語・ベトナム語のテキストペアを抽出

    Args:
        file_path: PPTXファイルのパス

    Returns:
        TextPairのリスト
    """
    file_path = Path(file_path)
    if not file_path.exists():
        raise FileNotFoundError(f"ファイルが見つかりません: {file_path}")

    prs = Presentation(str(file_path))
    pairs = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = extract_text_from_slide(slide)
        ja_texts, vi_texts = separate_by_language(texts)

        # テキストをペアリング（順序に基づく単純なマッチング）
        # 実際の資料の構造に合わせて調整が必要な場合あり
        for i, (ja, vi) in enumerate(zip(ja_texts, vi_texts)):
            pairs.append(TextPair(
                japanese=ja,
                vietnamese=vi,
                source_file=file_path.name,
                slide_number=slide_num
            ))

        # ペアにならなかった余りのテキストも記録（後で手動確認用）
        if len(ja_texts) != len(vi_texts):
            # 長い方の余りをログに記録
            pass

    return pairs


def extract_all_text_from_pptx(file_path: str | Path) -> dict:
    """
    PowerPointファイルから全テキストを言語別に抽出

    Args:
        file_path: PPTXファイルのパス

    Returns:
        {"japanese": [...], "vietnamese": [...], "other": [...]}
    """
    file_path = Path(file_path)
    if not file_path.exists():
        raise FileNotFoundError(f"ファイルが見つかりません: {file_path}")

    prs = Presentation(str(file_path))
    result = {"japanese": [], "vietnamese": [], "other": []}

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = extract_text_from_slide(slide)
        for text in texts:
            lang = detect_language(text)
            entry = {
                "text": text,
                "slide": slide_num,
                "language": lang
            }
            if lang == "ja":
                result["japanese"].append(entry)
            elif lang == "vi":
                result["vietnamese"].append(entry)
            else:
                result["other"].append(entry)

    return result
